#!/usr/bin/env python3
"""
Convert PDF pages to PPTX slides using pdf2image and python-pptx.

Usage:
    python pdf_to_pptx.py <input_pdf_path> <output_pptx_path>
"""

import io
import sys

try:
    from pdf2image import convert_from_path
except ImportError:
    sys.stderr.write(
        "pdf2image 모듈을 찾을 수 없습니다. `pip install pdf2image`로 설치하세요.\n"
    )
    sys.exit(2)

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.stderr.write(
        "python-pptx 모듈을 찾을 수 없습니다. `pip install python-pptx`로 설치하세요.\n"
    )
    sys.exit(2)


def px_to_emu(px: int, dpi: int) -> int:
    # 1 inch = 914400 EMU
    return int(px / dpi * 914400)


def main() -> int:
    if len(sys.argv) != 3:
        sys.stderr.write("Usage: pdf_to_pptx.py <input_pdf> <output_pptx>\n")
        return 1

    input_pdf, output_pptx = sys.argv[1:3]

    dpi = 200

    try:
        images = convert_from_path(input_pdf, dpi=dpi, fmt="png")
    except Exception as exc:  # pylint: disable=broad-except
        sys.stderr.write(f"PDF 페이지를 이미지로 변환하는 중 오류가 발생했습니다: {exc}\n")
        return 3

    if not images:
        prs = Presentation()
        prs.save(output_pptx)
        return 0

    prs = Presentation()

    first_image = images[0]
    prs.slide_width = Emu(px_to_emu(first_image.width, dpi))
    prs.slide_height = Emu(px_to_emu(first_image.height, dpi))

    blank_layout = prs.slide_layouts[6]

    for image in images:
        slide = prs.slides.add_slide(blank_layout)
        image_stream = io.BytesIO()
        image.save(image_stream, format="PNG")
        image_stream.seek(0)
        slide.shapes.add_picture(
            image_stream,
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    try:
        prs.save(output_pptx)
    except Exception as exc:  # pylint: disable=broad-except
        sys.stderr.write(f"PPTX 저장에 실패했습니다: {exc}\n")
        return 4

    return 0


if __name__ == "__main__":
    sys.exit(main())
